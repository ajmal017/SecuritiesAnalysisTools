import os

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # pylint: disable=no-name-in-module

from .slide_utils import slide_title_header


# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8

TEMP_DIR = os.path.join("output", "temp")


def make_MCI_slides(prs, analysis: dict):
    """Make MCI Slide

    Arguments:
        prs {pptx-object} -- presentation object
        analysis {dict} -- data object

    Returns:
        pptx-object -- modified pptx object
    """
    content = os.path.join(TEMP_DIR, "MCI.png")
    if os.path.exists(content):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, 'Market Composite Index')

        left = Inches(1.42)
        top = Inches(1.27)
        height = Inches(6.1)
        width = Inches(10.5)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    content = os.path.join(TEMP_DIR, "MCI_correlations.png")
    if os.path.exists(content):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, 'Market Composite Index')

        left = Inches(5.75)
        top = Inches(1.1)
        height = Inches(6.0)
        width = Inches(7.6)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

        # Add table here!
        if 'mci' in analysis.keys():
            num_rows = len(list(analysis['mci']['correlations'].keys())) + 2
            fund_key = list(analysis['mci']['correlations'].keys())[0]
            time_periods = [analysis['mci']['correlations'][fund_key][0]['period'],
                            analysis['mci']['correlations'][fund_key][1]['period']]

            # list of look back periods, having both B & R, plus name
            num_cols = 5  # len(analysis['MCI'][temp_key]) * 2 + 1

            left_loc = Inches(0.1)
            top_loc = Inches(1.1)
            table_width = Inches(5.75)
            table_height = Inches(6)

            table_placeholder = slide.shapes.add_table(num_rows,
                                                       num_cols,
                                                       left_loc,
                                                       top_loc,
                                                       table_width,
                                                       table_height)
            table = table_placeholder.table

            cell_1 = table.cell(0, 1)
            cell_2 = table.cell(0, 2)
            cell_1.merge(cell_2)
            cell_3 = table.cell(0, 3)
            cell_4 = table.cell(0, 4)
            cell_3.merge(cell_4)

            table.cell(1, 0).text = 'Fund'
            table.cell(0, 1).text = f"{time_periods[0]} Periods"
            table.cell(0, 3).text = f"{time_periods[1]} Periods"
            table.cell(1, 1).text = 'Beta'
            table.cell(1, 3).text = 'Beta'
            table.cell(1, 2).text = 'R-Squared'
            table.cell(1, 4).text = 'R-Squared'

            for i in range(5):
                table.cell(1, i).text_frame.paragraphs[0].font.size = Pt(15)
                table.cell(1, i).text_frame.paragraphs[0].font.bold = True

            for i, fund in enumerate(analysis['mci']['correlations'].keys()):
                table.cell(i+2, 0).text = fund
                table.cell(
                    i+2, 1).text = str(analysis['mci']['correlations'][fund][0]['beta'])
                table.cell(
                    i+2, 2).text = str(analysis['mci']['correlations'][fund][0]['r_squared'])
                table.cell(
                    i+2, 3).text = str(analysis['mci']['correlations'][fund][1]['beta'])
                table.cell(
                    i+2, 4).text = str(analysis['mci']['correlations'][fund][1]['r_squared'])

                table.cell(i+2, 0).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+2, 1).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+2, 2).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+2, 3).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+2, 4).text_frame.paragraphs[0].font.size = Pt(14)

    content = os.path.join(TEMP_DIR, "MCI_net_correlation.png")
    if os.path.exists(content):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, 'Market Composite Index')

        left = Inches(1.42)
        top = Inches(1.27)
        height = Inches(6.1)
        width = Inches(10.5)

        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    return prs


def make_BCI_slides(prs):
    """Make BCI Slides

    Arguments:
        prs {pptx-object} -- presentation object

    Returns:
        pptx-object -- presentation object
    """
    NUM_BOND_INDEXES = 3
    for i in range(NUM_BOND_INDEXES):
        if i == 0:
            filekey = 'Treasury'
        elif i == 1:
            filekey = 'Corporate'
        elif i == 2:
            filekey = 'International'
        else:
            return prs

        content = os.path.join(TEMP_DIR, f"{filekey}_BCI.png")
        if os.path.exists(content):

            title = f"{filekey} Bond Composite Index"
            slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
            slide = slide_title_header(slide, title)

            left = Inches(1.42)
            top = Inches(1.27)
            height = Inches(6.1)
            width = Inches(10.5)
            slide.shapes.add_picture(
                content, left, top, height=height, width=width)

    content = os.path.join(TEMP_DIR, "combined_BCI.png")
    if os.path.exists(content):

        title = f"Combined Bond Composite Indexes"
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, title)

        left = Inches(1.42)
        top = Inches(1.27)
        height = Inches(6.1)
        width = Inches(10.5)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    return prs


def make_CCI_slides(prs):
    """Make CCI Slides

    Arguments:
        prs {pptx-object} -- presentation object

    Returns:
        pptx-object -- presentation object
    """
    content = os.path.join(TEMP_DIR, "CCI_net_correlation.png")
    if os.path.exists(content):

        title = f"Correlation Composite Index"
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, title)

        left = Inches(1.42)
        top = Inches(1.27)
        height = Inches(6.1)
        width = Inches(10.5)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    return prs


def make_TCI_slides(prs):
    """Make TCI Slides

    Arguments:
        prs {pptx-object} -- pptx presentation object

    Returns:
        pptx-object -- pptx presentation object
    """
    content = os.path.join(TEMP_DIR, "tci.png")
    if os.path.exists(content):

        title = f"Type Composite Index"
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, title)

        left = Inches(1.42)
        top = Inches(1.27)
        height = Inches(6.1)
        width = Inches(10.5)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    return prs
